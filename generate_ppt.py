import sys, subprocess

# ensure python-pptx is installed
try:
    from pptx import Presentation
    from pptx.util import Pt, Inches
except Exception:
    subprocess.check_call([sys.executable, '-m', 'pip', 'install', 'python-pptx'])
    from pptx import Presentation
    from pptx.util import Pt, Inches

md_path = sys.argv[1]
output_path = sys.argv[2] if len(sys.argv) > 2 else md_path.rsplit('.',1)[0] + '.pptx'

with open(md_path, 'r', encoding='utf-8') as f:
    text = f.read()

# split slides on lines that are exactly '---'
parts = []
current = []
for line in text.splitlines():
    if line.strip() == '---':
        if current:
            parts.append('\n'.join(current).strip())
            current = []
    else:
        current.append(line)
if current:
    parts.append('\n'.join(current).strip())

prs = Presentation()
# define title slide from first H1 if present
if parts:
    first = parts[0]
    first_lines = [l.strip() for l in first.splitlines() if l.strip()]
    title_text = ''
    subtitle_text = ''
    if first_lines:
        for l in first_lines:
            if l.startswith('# '):
                title_text = l.lstrip('# ').strip()
                break
        if not title_text:
            title_text = first_lines[0].lstrip('#').strip()
        if len(first_lines) > 1:
            subtitle_text = ' '.join(first_lines[1:])
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = title_text or 'Healthcare Data for AI'
    subtitle.text = subtitle_text

from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT

def add_content_slide(prs, part_text):
    lines = [l.rstrip() for l in part_text.splitlines() if l.strip()]
    title = ''
    body_lines = []
    for idx, l in enumerate(lines):
        if l.startswith('### ') or l.startswith('## ') or l.startswith('# '):
            title = l.lstrip('# ').strip()
            body_lines = lines[idx+1:]
            break
    if not title and lines:
        title = lines[0]
        body_lines = lines[1:]
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    body = slide.shapes.placeholders[1].text_frame
    body.clear()
    for l in body_lines:
        if l.startswith('- '):
            p = body.add_paragraph()
            p.text = l.lstrip('- ').strip()
            p.level = 0
        elif l.startswith('|') and '|' in l:
            p = body.add_paragraph()
            p.text = l.strip()
            p.level = 0
        else:
            p = body.add_paragraph()
            p.text = l.strip()
            p.level = 0

start_idx = 1 if parts and parts[0].strip() else 0
for part in parts[start_idx:]:
    add_content_slide(prs, part)

prs.save(output_path)
print(f'PPTX written to: {output_path}')
